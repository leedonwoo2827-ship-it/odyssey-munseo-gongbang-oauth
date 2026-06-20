"""McKinsey식 덱 PPTX 생성 — 템플릿 슬라이드(A/B) 복제 + 네이티브 차트.

일반 pptx_gen(제목+불릿)과 달리, LLM 이 낸 구조화된 deck spec(JSON)을 받아
템플릿의 박스 안에 네이티브 PowerPoint 차트/표/KPI 를 렌더링한다.

payload(deck spec) 스키마:
  {
    "title": "...", "subtitle": "...",            # (선택) 표지용 메타
    "slides": [
      {
        "template": "A" | "B",                    # A=박스1, B=박스2
        "chapter": "01  도입", "title": "...", "subtitle": "...",
        "boxes": [
          {
            "header": "...", "unit": "단위: %",
            "kind": "column|bar|line|doughnut|table|kpi",
            "categories": [...], "series": [{"name":"..","values":[..]}],   # column/bar/line
            "labels": [...], "values": [...],                               # doughnut
            "headers": [...], "rows": [[...]],                              # table
            "kpis": [{"value":"34%","label":"..","desc":".."}],            # kpi
            "number_format": "0",
            "insight": "...", "source": "Source: .."
          }
        ]
      }
    ]
  }

템플릿(회사/영상 brandlogy)이 슬라이드 1·2 를 가지면 이를 복제(배경/박스/로고/푸터 보존).
템플릿이 없거나 슬라이드가 부족하면 둥근 사각형을 직접 그려 폴백한다.
"""
from __future__ import annotations

import copy
import json
import os
from typing import Any, Dict, List, Optional, Tuple

# ── 브랜딩 변수 (회사화 시 교체) ─────────────────────────────────────────
FONT = "Pretendard"
BLUE_L1 = (0x00, 0x19, 0xB5)
BLUE_L3 = (0xB8, 0xC8, 0xF5)
BLUE_L4 = (0x8A, 0x9D, 0xD8)
GRAY_25 = (0xBF, 0xBF, 0xBF)
GRAY_50 = (0x7F, 0x7F, 0x7F)
GRAY_75 = (0x40, 0x40, 0x40)
WHITE = (0xFF, 0xFF, 0xFF)
SERIES_PALETTE = [BLUE_L1, BLUE_L3, BLUE_L4, GRAY_25]

# ── 박스 지오메트리(EMU) — 템플릿 XML 실측값 ───────────────────────────
BOX_A = dict(x=488950, y=2060576, w=11212271, h=4076700)
BOX_B_LEFT = dict(x=488950, y=2060576, w=5462269, h=4076700)
BOX_B_RIGHT = dict(x=6240782, y=2060576, w=5462269, h=4076700)
PAD = 274320

PH_CHAPTER = "챕터명을 입력해 주세요"
PH_TITLE = "제목을 입력해 주세요"
PH_SUBTITLE = "부제목을 입력해 주세요"


# ── payload 정규화 ──────────────────────────────────────────────────────
def _coerce(payload: Any) -> Dict[str, Any]:
    if isinstance(payload, str):
        try:
            payload = json.loads(payload)
        except Exception:
            return {"title": "", "subtitle": "", "slides": []}
    if isinstance(payload, dict) and isinstance(payload.get("slides"), list):
        return payload
    return {"title": "", "subtitle": "", "slides": []}


def _norm_series(series: Any) -> List[Tuple[str, List]]:
    out: List[Tuple[str, List]] = []
    if isinstance(series, dict):
        for name, vals in series.items():
            out.append((str(name), list(vals)))
    elif isinstance(series, list):
        for s in series:
            if isinstance(s, dict):
                out.append((str(s.get("name", "")), list(s.get("values", []))))
            elif isinstance(s, (list, tuple)) and len(s) == 2:
                out.append((str(s[0]), list(s[1])))
    return out


# ── 메인 ────────────────────────────────────────────────────────────────
def render(payload: Any, out_path: str, template: Optional[str] = None) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    data = _coerce(payload)
    slides_spec = data.get("slides") or []

    if template and os.path.exists(template):
        prs = Presentation(template)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    base_slides = list(prs.slides)
    has_templates = len(base_slides) >= 2
    src_a = base_slides[0] if has_templates else None
    src_b = base_slides[1] if has_templates else None

    for spec in slides_spec:
        _render_slide(prs, src_a, src_b, spec)

    if not slides_spec:
        # 빈 결과 방어: 표지 한 장
        _render_slide(prs, src_a, src_b, {
            "template": "A", "chapter": "", "title": data.get("title", "내용"),
            "subtitle": data.get("subtitle", ""), "boxes": [],
        })

    if has_templates:
        _delete_front(prs, 2)

    prs.save(out_path)
    return out_path


def _render_slide(prs, src_a, src_b, spec: Dict[str, Any]):
    tmpl = str(spec.get("template", "B")).upper()
    if tmpl == "A":
        slide = _new_slide(prs, src_a, boxes=[BOX_A])
        boxes = [BOX_A]
    else:
        slide = _new_slide(prs, src_b, boxes=[BOX_B_LEFT, BOX_B_RIGHT])
        boxes = [BOX_B_LEFT, BOX_B_RIGHT]

    _set_meta(slide, PH_CHAPTER, spec.get("chapter", ""))
    _set_meta(slide, PH_TITLE, spec.get("title", ""))
    _set_meta(slide, PH_SUBTITLE, spec.get("subtitle", ""))

    for box, content in zip(boxes, spec.get("boxes", []) or []):
        try:
            _render_box(slide, box, content)
        except Exception as e:  # 한 박스 실패가 전체 덱을 막지 않게
            _add_text(slide, box["x"] + PAD, box["y"] + PAD, box["w"] - 2 * PAD, 400000,
                      f"[렌더 오류: {type(e).__name__}]", 10, GRAY_50)
    return slide


# ── 슬라이드 생성(복제 또는 직접 그리기) ────────────────────────────────
def _new_slide(prs, src, boxes):
    if src is not None:
        slide = _duplicate(prs, src)
    else:
        slide = _blank_with_boxes(prs, boxes)
    return slide


def _duplicate(prs, src_slide):
    layout = src_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)
    for shape in src_slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))
    return new_slide


def _blank_with_boxes(prs, boxes):
    """템플릿이 없을 때: 빈 레이아웃 + 둥근 사각형 + 메타 텍스트박스."""
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    for b in boxes:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Emu(b["x"]), Emu(b["y"]), Emu(b["w"]), Emu(b["h"]))
        shp.fill.solid(); _set_rgb(shp.fill.fore_color, WHITE)
        shp.line.color.rgb = _rgb(GRAY_25); shp.line.width = Emu(9525)
        shp.shadow.inherit = False
        if shp.has_text_frame:
            shp.text_frame.clear()
    # 메타 텍스트박스(템플릿 좌표와 동일)
    _add_text(slide, 385854, 279533, 3324826, 246221, PH_CHAPTER, 10, _BLUE_L2(), bold=True)
    _add_text(slide, 385854, 777736, 6000000, 461665, PH_TITLE, 24, (0x26, 0x26, 0x26), bold=True)
    _add_text(slide, 385854, 1345865, 6000000, 276999, PH_SUBTITLE, 12, GRAY_75)
    return slide


def _BLUE_L2():
    return (0x4F, 0x6F, 0xF0)


def _delete_front(prs, n):
    lst = prs.slides._sldIdLst
    for sldId in list(lst)[:n]:
        lst.remove(sldId)


# ── 텍스트/폰트 ─────────────────────────────────────────────────────────
def _rgb(t):
    from pptx.dml.color import RGBColor
    return RGBColor(*t)


def _set_rgb(color_fmt, t):
    color_fmt.rgb = _rgb(t)


def _set_run(run, name, size, color, bold=False):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _set_meta(slide, marker, new_text):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if marker in run.text:
                    run.text = new_text or ""
                    return True
    return False


def _add_text(slide, x, y, w, h, text, size, color, bold=False, align=None,
              anchor=None, font=FONT):
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor or MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align or PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text or ""
    _set_run(run, font, size, color, bold)
    return tb


# ── 박스 내부(헤더/단위/시각화/인사이트/출처) ──────────────────────────
def _render_box(slide, box, content):
    from pptx.enum.text import MSO_ANCHOR
    x, y, w, h = box["x"], box["y"], box["w"], box["h"]
    ix = x + PAD
    iw = w - 2 * PAD
    cur = y + int(PAD * 0.7)

    _add_text(slide, ix, cur, iw, 300000, content.get("header", ""), 13, GRAY_75, bold=True)
    cur += 330000
    if content.get("unit"):
        _add_text(slide, ix, cur, iw, 200000, content["unit"], 9, GRAY_50)
    cur += 230000

    src_h, ins_h = 170000, 260000
    src_y = y + h - int(PAD * 0.5) - src_h
    ins_y = src_y - ins_h
    viz_rect = (ix, cur, iw, ins_y - cur - 40000)

    kind = str(content.get("kind", "")).lower()
    if kind in ("column", "bar", "line"):
        _add_cat_chart(slide, viz_rect, kind, content.get("categories", []),
                       _norm_series(content.get("series")),
                       content.get("number_format", "0"))
    elif kind == "doughnut":
        _add_doughnut(slide, viz_rect, content.get("labels", []), content.get("values", []))
    elif kind == "table":
        _add_table(slide, viz_rect, content.get("headers", []), content.get("rows", []))
    elif kind == "kpi":
        _add_kpis(slide, viz_rect, content.get("kpis", []))

    _add_text(slide, ix, ins_y, iw, ins_h, content.get("insight", ""), 11, GRAY_75,
              bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, ix, src_y, iw, src_h, content.get("source", ""), 8, GRAY_50)


# ── 네이티브 차트/표/KPI ────────────────────────────────────────────────
def _chart_fonts(chart, size=9):
    from pptx.util import Pt
    try:
        chart.font.size = Pt(size); chart.font.name = FONT
    except Exception:
        pass


def _add_cat_chart(slide, rect, kind, categories, series, number_format):
    from pptx.util import Emu, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    ctype = {"column": XL_CHART_TYPE.COLUMN_CLUSTERED,
             "bar": XL_CHART_TYPE.BAR_CLUSTERED,
             "line": XL_CHART_TYPE.LINE_MARKERS}[kind]
    cd = CategoryChartData()
    cd.categories = categories or [""]
    if not series:
        series = [("값", [0] * len(cd.categories))]
    for name, vals in series:
        cd.add_series(name, vals)
    x, y, cx, cy = rect
    chart = slide.shapes.add_chart(ctype, Emu(x), Emu(y), Emu(cx), Emu(cy), cd).chart
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    _chart_fonts(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = number_format
    plot.data_labels.number_format_is_linked = False
    try:
        plot.data_labels.font.size = Pt(9); plot.data_labels.font.name = FONT
    except Exception:
        pass
    for i, s in enumerate(chart.series):
        try:
            s.format.fill.solid(); s.format.fill.fore_color.rgb = _rgb(SERIES_PALETTE[i % 4])
        except Exception:
            pass


def _add_doughnut(slide, rect, labels, values):
    from pptx.util import Emu
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    cd = CategoryChartData()
    cd.categories = labels or [""]
    cd.add_series("점유율", values or [0])
    x, y, cx, cy = rect
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Emu(x), Emu(y), Emu(cx), Emu(cy), cd).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    _chart_fonts(chart)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0"%"'
    plot.data_labels.number_format_is_linked = False
    try:
        plot.data_labels.position = XL_LABEL_POSITION.CENTER
    except Exception:
        pass
    for i, pt in enumerate(chart.series[0].points):
        try:
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = _rgb(SERIES_PALETTE[i % 4])
        except Exception:
            pass


def _add_table(slide, rect, headers, rows):
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    if not headers:
        return
    x, y, cx, cy = rect
    nrows, ncols = len(rows) + 1, len(headers)
    table = slide.shapes.add_table(nrows, ncols, Emu(x), Emu(y), Emu(cx), Emu(cy)).table
    for c, htxt in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid(); _set_rgb(cell.fill.fore_color, GRAY_75)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(htxt); _set_run(r, FONT, 11, WHITE, bold=True)
    for ri, row in enumerate(rows, start=1):
        for c in range(ncols):
            val = row[c] if c < len(row) else ""
            cell = table.cell(ri, c)
            cell.fill.solid(); _set_rgb(cell.fill.fore_color, WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            color = BLUE_L1 if c == ncols - 1 else GRAY_75
            _set_run(r, FONT, 10, color, bold=(c == ncols - 1))


def _add_kpis(slide, rect, kpis):
    x, y, w, h = rect
    n = max(1, len(kpis))
    bh = h // n
    for i, k in enumerate(kpis):
        by = y + i * bh
        _add_text(slide, x, by + 20000, w, 520000, str(k.get("value", "")), 30, BLUE_L1, bold=True)
        _add_text(slide, x, by + 560000, w, 240000, str(k.get("label", "")), 12, GRAY_75, bold=True)
        if k.get("desc"):
            _add_text(slide, x, by + 800000, w, 220000, str(k["desc"]), 9, GRAY_50)
