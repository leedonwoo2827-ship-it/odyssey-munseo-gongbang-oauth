"""디자인 덱 PPTX 생성 — 도형 조립형(카드/타임라인/진행바/KPI/간지).

mckinsey_pptx_gen(고정 박스+네이티브 차트)·pptx_gen(제목+불릿)과 달리, 슬라이드마다
'블록'을 도형으로 조립해 리디자인 수준의 발표자료를 만든다. 회사 템플릿(pptx_template)의
마스터(로고/브랜딩)를 그대로 쓰고, 슬라이드 위에 직접 도형을 그린다.

deck spec(JSON):
  {"slides": [
    {"type":"cover","title":..,"subtitle":..},
    {"type":"toc","chapter":"OVERVIEW","title":"목차","items":[{"no":"01","part":"Part I","title":".."}]},
    {"type":"divider","no":"01","part":"Part I","title":".."},
    {"type":"content","chapter":"PART I · ..","title":"결론 문장","subtitle":"한 줄",
     "blocks":[
        {"kind":"progress","header":"단계별 완료율","items":[{"label":"2024.10.31","sub":"12/17","pct":71}]},
        {"kind":"kpi_cards","cards":[{"header":"프로젝트 종료식","value":"2025.8.11","desc":".."}]},
        {"kind":"timeline","header":"..","items":[{"date":"2025.4.8","text":".."}]},
        {"kind":"bullets","header":"..","items":["..",".."]},
        {"kind":"table","header":"..","headers":["a","b"],"rows":[["1","2"]]}
     ]}
  ]}
값에 '확인 필요'/'?'가 들어가면 앰버 강조(작성자가 채울 자리).
"""
from __future__ import annotations

import os
from typing import Any, Dict, List, Optional

# ── 디자인 토큰(리디자인 덱에서 추출) ────────────────────────────────────
FONT = "Pretendard"
PRIMARY = "4F6EF1"   # 인디고(주색)
NAVY = "0E2841"      # 제목
BODY = "333333"      # 본문
SUB = "6B7280"       # 보조 텍스트
CARD_BG = "F4F6FB"   # 카드 배경
BORDER = "DCE1EA"    # 카드 테두리
TEAL = "00A88E"      # 진행/긍정
TRACK = "E3E8EF"     # 진행바 트랙
AMBER = "D98613"     # 확인 필요(강조)
AMBER_BG = "FBF1DE"  # 확인 필요 배경
AMBER_TX = "9A5E08"  # 확인 필요 텍스트
WHITE = "FFFFFF"

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)
MX = int(0.92 * EMU)          # 좌우 여백
CONTENT_W = SW - 2 * MX


def _need(s: Any) -> bool:
    t = str(s or "")
    return ("확인 필요" in t) or t.strip() in ("?", "?%", "? %", "[확인 필요]")


# ── 저수준 도형 헬퍼 ─────────────────────────────────────────────────────
def _rgb(hexs: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hexs)


def _strip_slides(prs) -> None:
    from pptx.oxml.ns import qn
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        rId = sldId.get(qn("r:id"))
        lst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass


def _blank(prs):
    """마스터(로고)는 유지하되 레이아웃 placeholder 는 제거한 빈 슬라이드."""
    layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]
    s = prs.slides.add_slide(layout)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    return s


def _box(slide, x, y, w, h, *, fill=None, line=None, line_w=1.0, round_=True):
    from pptx.util import Emu, Pt
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = _rgb(line); shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, *, align="l", anchor="t"):
    """runs: [(text, size, color, bold)] 또는 [[run,...] per paragraph]."""
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}[anchor]
    paras = runs if runs and isinstance(runs[0], list) else [runs]
    al = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.color.rgb = _rgb(color); r.font.bold = bold
    return tb


def _header(slide, chapter, title, subtitle=""):
    """슬라이드 상단: 챕터(소형 블루) + 제목(네이비 26) + 부제(보조)."""
    y = int(0.82 * EMU)
    if chapter:
        _text(slide, MX, y, CONTENT_W, int(0.3 * EMU), [(chapter, 11.5, PRIMARY, True)])
        y += int(0.34 * EMU)
    _text(slide, MX, y, CONTENT_W, int(0.7 * EMU), [(title, 26, NAVY, True)])
    y += int(0.78 * EMU)
    if subtitle:
        _text(slide, MX, y, CONTENT_W, int(0.4 * EMU), [(subtitle, 13.5, SUB, False)])
        y += int(0.5 * EMU)
    return y + int(0.15 * EMU)


# ── 블록 렌더러 ──────────────────────────────────────────────────────────
def _block_header(slide, x, y, w, header, unit=""):
    if header:
        _text(slide, x, y, w, int(0.32 * EMU), [(header, 14.5, NAVY, True)])
    if unit:
        _text(slide, x, y + int(0.34 * EMU), w, int(0.22 * EMU), [(unit, 10.5, SUB, False)])
    return y + int((0.62 if unit else 0.42) * EMU)


def _progress(slide, x, y, w, h, blk):
    cy = _block_header(slide, x, y, w, blk.get("header", ""), blk.get("unit", ""))
    items = blk.get("items", [])[:4] or [{}]
    row_h = min(int(1.0 * EMU), (y + h - cy) // max(1, len(items)))
    bar_x = x + int(2.0 * EMU)
    bar_w = w - int(2.0 * EMU) - int(1.1 * EMU)
    for it in items:
        pct = it.get("pct")
        need = pct is None or _need(it.get("pct"))
        # 라벨 + 보조
        _text(slide, x, cy + int(0.06 * EMU), int(1.9 * EMU), int(0.5 * EMU),
              [[(str(it.get("label", "")), 13, NAVY, True)],
               [(str(it.get("sub", "")), 10.5, SUB, False)]])
        by = cy + int(0.12 * EMU); bh = int(0.36 * EMU)
        if need:
            box = _box(slide, bar_x, by, bar_w, bh, fill=None, line=AMBER, line_w=1.25)
            try:
                box.line.dash_style = None
            except Exception:
                pass
            _text(slide, bar_x, by, bar_w, bh, [("확인 필요", 11, AMBER_TX, False)],
                  align="c", anchor="m")
            _text(slide, bar_x + bar_w + int(0.12 * EMU), cy, int(1.0 * EMU), int(0.55 * EMU),
                  [("?", 24, AMBER, True)], anchor="m")
        else:
            p = max(0, min(100, float(pct)))
            _box(slide, bar_x, by, bar_w, bh, fill=TRACK, line=None)
            _box(slide, bar_x, by, int(bar_w * p / 100), bh, fill=TEAL, line=None)
            _text(slide, bar_x + bar_w + int(0.12 * EMU), cy, int(1.0 * EMU), int(0.55 * EMU),
                  [(f"{int(p)}%", 24, TEAL, True)], anchor="m")
        cy += row_h


def _kpi_cards(slide, x, y, w, h, blk):
    cards = blk.get("cards", [])[:3] or [{}]
    gap = int(0.18 * EMU)
    ch = (h - gap * (len(cards) - 1)) // max(1, len(cards))
    cy = y
    for c in cards:
        need = _need(c.get("value"))
        bg = AMBER_BG if need else CARD_BG
        _box(slide, x, cy, w, ch, fill=bg, line=None)
        pad = int(0.22 * EMU)
        iw = w - 2 * pad
        # header → value → desc 를 위에서부터 쌓아 겹치지 않게.
        _text(slide, x + pad, cy + pad, iw, int(0.28 * EMU),
              [(str(c.get("header", "")), 12.5, (AMBER_TX if need else PRIMARY), True)])
        _text(slide, x + pad, cy + pad + int(0.3 * EMU), iw, int(0.5 * EMU),
              [(str(c.get("value", "")), 25, (AMBER if need else NAVY), True)])
        if c.get("desc"):
            _text(slide, x + pad, cy + pad + int(0.84 * EMU), iw, ch - int(1.06 * EMU),
                  [(str(c["desc"]), 10.5, (AMBER_TX if need else SUB), False)])
        cy += ch + gap


def _timeline(slide, x, y, w, h, blk):
    cy = _block_header(slide, x, y, w, blk.get("header", ""))
    items = blk.get("items", [])[:6] or [{}]
    line_x = x + int(2.0 * EMU)
    row_h = min(int(0.95 * EMU), (y + h - cy) // max(1, len(items)))
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE
    # 세로선
    if len(items) > 1:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(line_x - int(0.01 * EMU)),
                                    Emu(cy + row_h // 2), Emu(int(0.02 * EMU)),
                                    Emu(row_h * (len(items) - 1)))
        ln.fill.solid(); ln.fill.fore_color.rgb = _rgb(BORDER); ln.line.fill.background()
        ln.shadow.inherit = False
    for it in items:
        d = int(0.16 * EMU)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(line_x - d // 2),
                                     Emu(cy + row_h // 2 - d // 2), Emu(d), Emu(d))
        dot.fill.solid(); dot.fill.fore_color.rgb = _rgb(PRIMARY); dot.line.fill.background()
        dot.shadow.inherit = False
        _text(slide, x, cy, int(1.85 * EMU) - int(0.12 * EMU), row_h,
              [(str(it.get("date", "")), 12.5, PRIMARY, True)], align="r", anchor="m")
        txt = str(it.get("text", ""))
        col = AMBER_TX if _need(txt) else BODY
        _text(slide, line_x + int(0.22 * EMU), cy, w - (line_x - x) - int(0.3 * EMU), row_h,
              [(txt, 14, col, False)], anchor="m")
        cy += row_h


def _bullets(slide, x, y, w, h, blk):
    cy = _block_header(slide, x, y, w, blk.get("header", ""))
    for it in blk.get("items", [])[:8]:
        col = AMBER_TX if _need(it) else BODY
        _text(slide, x, cy, w, int(0.42 * EMU),
              [[("•  ", 14, PRIMARY, True), (str(it), 14, col, False)]])
        cy += int(0.46 * EMU)


def _table(slide, x, y, w, h, blk):
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    cy = _block_header(slide, x, y, w, blk.get("header", ""), blk.get("unit", ""))
    headers = blk.get("headers", []) or []
    rows = blk.get("rows", []) or []
    if not headers and not rows:
        return
    ncol = len(headers) or (len(rows[0]) if rows else 1)
    nrow = len(rows) + 1
    th = min(int(0.55 * EMU), max(int(0.42 * EMU), (y + h - cy) // max(1, nrow)))
    gt = slide.shapes.add_table(nrow, ncol, Emu(x), Emu(cy), Emu(w),
                                Emu(th * nrow)).table
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = _rgb(NAVY)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(htxt)
        r.font.name = FONT; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = _rgb(WHITE)
    for i, row in enumerate(rows, start=1):
        for j in range(ncol):
            val = str(row[j]) if j < len(row) else ""
            c = gt.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = _rgb(WHITE if i % 2 else CARD_BG)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.name = FONT; r.font.size = Pt(11)
            r.font.color.rgb = _rgb(AMBER_TX if _need(val) else BODY)


_BLOCKS = {"progress": _progress, "kpi_cards": _kpi_cards, "timeline": _timeline,
           "bullets": _bullets, "table": _table}


# ── 슬라이드 타입 ────────────────────────────────────────────────────────
def _slide_cover(prs, spec):
    s = _blank(prs)
    _text(s, MX, int(3.0 * EMU), CONTENT_W, int(1.0 * EMU),
          [(str(spec.get("title", "")), 36, NAVY, True)], align="c", anchor="m")
    if spec.get("subtitle"):
        _text(s, MX, int(4.1 * EMU), CONTENT_W, int(0.6 * EMU),
              [(str(spec["subtitle"]), 15, SUB, False)], align="c")


def _slide_toc(prs, spec):
    s = _blank(prs)
    _header(s, spec.get("chapter", "OVERVIEW"), spec.get("title", "목차"))
    items = spec.get("items", [])[:6]
    gap = int(0.25 * EMU)
    cw = (CONTENT_W - gap) // 2
    chh = int(0.95 * EMU)
    y0 = int(2.5 * EMU)
    for i, it in enumerate(items):
        col, row = i % 2, i // 2
        cx = MX + col * (cw + gap)
        cy = y0 + row * (chh + gap)
        _box(s, cx, cy, cw, chh, fill=CARD_BG, line=None)
        # 번호 원
        from pptx.util import Emu
        from pptx.enum.shapes import MSO_SHAPE
        d = int(0.55 * EMU)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx + int(0.22 * EMU)),
                                  Emu(cy + (chh - d) // 2), Emu(d), Emu(d))
        circ.fill.solid(); circ.fill.fore_color.rgb = _rgb(PRIMARY); circ.line.fill.background()
        circ.shadow.inherit = False
        _text(s, cx + int(0.22 * EMU), cy + (chh - d) // 2, d, d,
              [(str(it.get("no", "")), 14, WHITE, True)], align="c", anchor="m")
        tx = cx + int(1.0 * EMU)
        _text(s, tx, cy + int(0.2 * EMU), cw - int(1.2 * EMU), int(0.3 * EMU),
              [(str(it.get("part", "")), 11.5, PRIMARY, True)])
        _text(s, tx, cy + int(0.5 * EMU), cw - int(1.2 * EMU), int(0.35 * EMU),
              [(str(it.get("title", "")), 15, NAVY, True)])


def _slide_divider(prs, spec):
    s = _blank(prs)
    _text(s, MX, int(1.6 * EMU), int(5 * EMU), int(3 * EMU),
          [(str(spec.get("no", "")), 210, "EAEEFC", True)])
    _text(s, MX, int(3.7 * EMU), CONTENT_W, int(0.4 * EMU),
          [(str(spec.get("part", "")), 14, PRIMARY, True)])
    _text(s, MX, int(4.15 * EMU), CONTENT_W, int(0.8 * EMU),
          [(str(spec.get("title", "")), 30, NAVY, True)])


def _slide_content(prs, spec):
    s = _blank(prs)
    top = _header(s, spec.get("chapter", ""), spec.get("title", ""), spec.get("subtitle", ""))
    blocks = spec.get("blocks", []) or []
    area_y, area_h = top, int(6.95 * EMU) - top
    side = [b for b in blocks if b.get("kind") == "kpi_cards"]
    main = [b for b in blocks if b.get("kind") != "kpi_cards"]
    if side and main:
        mw = int(CONTENT_W * 0.56)
        sx = MX + mw + int(0.3 * EMU)
        sw = SW - MX - sx
        _render_main(s, MX, area_y, mw, area_h, main)
        # 사이드 kpi 카드들 합쳐서
        allcards = [c for b in side for c in b.get("cards", [])]
        _kpi_cards(s, sx, area_y, sw, area_h, {"cards": allcards})
    else:
        _render_main(s, MX, area_y, CONTENT_W, area_h, blocks)


def _render_main(slide, x, y, w, h, blocks):
    blocks = blocks or []
    if not blocks:
        return
    bh = (h - int(0.2 * EMU) * (len(blocks) - 1)) // len(blocks)
    cy = y
    for b in blocks:
        fn = _BLOCKS.get(b.get("kind"))
        if fn:
            fn(slide, x, cy, w, bh, b)
        cy += bh + int(0.2 * EMU)


_TYPES = {"cover": _slide_cover, "toc": _slide_toc, "divider": _slide_divider,
          "content": _slide_content}


def render(payload: Any, out_path: str, template: Optional[str] = None) -> str:
    from pptx import Presentation
    from pptx.util import Inches
    data = payload if isinstance(payload, dict) else {}
    slides = data.get("slides") or []
    if template and os.path.exists(template):
        prs = Presentation(template)
        _strip_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    for spec in slides:
        fn = _TYPES.get(str(spec.get("type", "content")).lower(), _slide_content)
        try:
            fn(prs, spec)
        except Exception:
            _slide_content(prs, {"chapter": spec.get("chapter", ""),
                                 "title": spec.get("title", ""), "blocks": []})
    prs.save(out_path)
    return out_path
