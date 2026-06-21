"""PPTX 생성 — 어떤 템플릿이든 견고하게 표지/본문 슬라이드를 채운다.

설계 원칙(중요 — 깨진 템플릿 대응):
  - placeholder 를 '인덱스'가 아니라 '유형(type)'으로 찾는다. 회사 양식처럼 idx 가
    105/106/108 식으로 비표준이어도 동작.
  - 텍스트는 word_wrap + auto_size(TEXT_TO_FIT_SHAPE) + 명시 폰트 상한으로 폭주 방지.
  - 우리가 채우지 않은 본문/제목 placeholder 의 '안내문구'(예: 'PPT Master Template')는
    빈 값으로 지워 화면에 남지 않게 한다.
  - 템플릿에 이미 들어있는 샘플 슬라이드는 제거하고 우리 슬라이드만 추가한다.

payload(slides 모드):
  {"title": "...", "subtitle": "...",
   "slides": [{"title": "...", "bullets": ["...", ...]}, ...]}
Markdown 이 와도 _coerce_slides 가 슬라이드 구조로 변환한다.
"""
from __future__ import annotations

import os
from typing import Any, Dict, List, Optional

from . import md_gen

# 폰트 상한(pt) — 템플릿 기본값이 거대해도 이 값으로 덮어 폭주 방지
_TITLE_PT = 28
_SUBTITLE_PT = 18
_BODY_PT = 16
_COVER_TITLE_PT = 30


def _coerce_slides(payload: Any) -> Dict:
    """payload 가 Markdown 이면 슬라이드 구조로 최대한 변환."""
    if isinstance(payload, dict) and "slides" in payload:
        return payload
    from .. import mdblocks
    md = payload if isinstance(payload, str) else md_gen.to_markdown(payload)
    blocks = mdblocks.parse(md)
    title, subtitle = "", ""
    slides: List[Dict] = []
    cur: Optional[Dict] = None
    for b in blocks:
        if b["type"] == "heading" and b["level"] == 1 and not title:
            title = b["text"]
        elif b["type"] == "heading":
            cur = {"title": b["text"], "bullets": []}
            slides.append(cur)
        elif b["type"] == "bullet":
            (cur or _new(slides))["bullets"].append(b["text"])
        elif b["type"] == "para":
            if cur is None and not subtitle:
                subtitle = b["text"]
            else:
                (cur or _new(slides))["bullets"].append(b["text"])
    if not slides:
        slides = [{"title": title or "내용", "bullets": [subtitle] if subtitle else []}]
    return {"title": title or "제목", "subtitle": subtitle, "slides": slides}


def _new(slides: List[Dict]) -> Dict:
    d = {"title": "내용", "bullets": []}
    slides.append(d)
    return d


# ── placeholder 유형 분류 ────────────────────────────────────────────
def _ph_types():
    from pptx.enum.shapes import PP_PLACEHOLDER as PP
    title = {PP.TITLE, PP.CENTER_TITLE, PP.VERTICAL_TITLE}
    subtitle = {PP.SUBTITLE}
    body = {PP.BODY, PP.OBJECT, PP.VERTICAL_BODY}
    skip = {PP.SLIDE_NUMBER, PP.DATE, PP.FOOTER, PP.HEADER, PP.PICTURE,
            PP.BITMAP, PP.MEDIA_CLIP, PP.TABLE, PP.CHART, PP.ORG_CHART}
    return title, subtitle, body, skip


def _is_type(ph, types) -> bool:
    try:
        return ph.placeholder_format.type in types
    except Exception:
        return False


def _area(ph) -> int:
    try:
        return int(ph.width) * int(ph.height)
    except Exception:
        return 0


def _pick_layouts(prs):
    """표지/본문 레이아웃을 유형·이름으로 견고하게 고른다."""
    T, S, B, _ = _ph_types()
    layouts = list(prs.slide_layouts)

    def has(lay, types):
        return any(_is_type(p, types) for p in lay.placeholders)

    def name_has(lay, words):
        n = (lay.name or "").lower()
        return any(w in n for w in words)

    cover = None
    for lay in layouts:  # 1순위: 이름이 표지/제목 슬라이드
        if name_has(lay, ("표지", "제목 슬라이드", "title slide", "cover")) and has(lay, T):
            cover = lay
            break
    if cover is None:
        for lay in layouts:  # 2순위: 제목 + 부제
            if has(lay, T) and has(lay, S):
                cover = lay
                break
    if cover is None:
        for lay in layouts:
            if has(lay, T):
                cover = lay
                break
    cover = cover or layouts[0]

    body = None
    for lay in layouts:  # 1순위: 이름이 본문/제목 및 내용
        if name_has(lay, ("기본본문", "본문", "제목 및 내용", "title and content", "content")) \
                and has(lay, T) and has(lay, B) and lay is not cover:
            body = lay
            break
    if body is None:
        for lay in layouts:  # 2순위: 제목 + 본문 placeholder
            if has(lay, T) and has(lay, B) and lay is not cover:
                body = lay
                break
    if body is None:
        for lay in layouts:
            if has(lay, B) and lay is not cover:
                body = lay
                break
    body = body or cover
    return cover, body


def _strip_slides(prs) -> None:
    """템플릿에 들어있던 샘플 슬라이드를 모두 제거(파트+관계까지).

    sldIdLst 항목만 지우면 슬라이드 '파트'(slide1.xml 등)가 패키지에 남아, 새 슬라이드가
    같은 이름(slide1.xml)으로 추가될 때 '중복 파트명'으로 파일이 손상된다. 관계(rel)를
    함께 끊어 옛 파트가 저장되지 않게 한다.
    """
    from pptx.oxml.ns import qn
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        sldIdLst.remove(sldId)
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass


def render(payload: Any, out_path: str, template: Optional[str] = None) -> str:
    from pptx import Presentation
    from pptx.util import Inches

    data = _coerce_slides(payload)

    if template and os.path.exists(template):
        prs = Presentation(template)
        _strip_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    cover_layout, body_layout = _pick_layouts(prs)

    # ── 표지 ──
    s = prs.slides.add_slide(cover_layout)
    _fill_cover(s, data.get("title", ""), data.get("subtitle", ""))

    # ── 본문 슬라이드 ──
    for slide in data.get("slides", []):
        s = prs.slides.add_slide(body_layout)
        _fill_body(s, slide.get("title", ""), slide.get("bullets", []))

    prs.save(out_path)
    return out_path


# ── 채우기 헬퍼 ──────────────────────────────────────────────────────
def _set_text(ph, lines: List[str], size_pt: int, bold: bool = False) -> None:
    from pptx.util import Pt
    from pptx.enum.text import MSO_AUTO_SIZE
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    lines = lines or [""]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(ln)
        for run in p.runs:
            run.font.size = Pt(size_pt)
            if bold:
                run.font.bold = True
        if not p.runs:  # 빈 줄도 폰트 지정
            p.font.size = Pt(size_pt)


def _clear(ph) -> None:
    try:
        if ph.has_text_frame:
            ph.text_frame.clear()
    except Exception:
        pass


def _fill_cover(slide, title: str, subtitle: str) -> None:
    T, S, B, SKIP = _ph_types()
    used = set()
    # 제목
    title_ph = next((p for p in slide.placeholders if _is_type(p, T)), None)
    if title_ph is None:
        title_ph = next((p for p in slide.placeholders), None)
    if title_ph is not None:
        _set_text(title_ph, _wrap(title), _COVER_TITLE_PT, bold=True)
        used.add(title_ph.placeholder_format.idx)
    # 부제: SUBTITLE 우선, 없으면 가장 큰 BODY
    sub_ph = next((p for p in slide.placeholders if _is_type(p, S)), None)
    if sub_ph is None:
        bodies = [p for p in slide.placeholders
                  if _is_type(p, B) and p.placeholder_format.idx not in used]
        sub_ph = max(bodies, key=_area) if bodies else None
    if sub_ph is not None and subtitle:
        _set_text(sub_ph, _wrap(subtitle), _SUBTITLE_PT)
        used.add(sub_ph.placeholder_format.idx)
    _clear_unused(slide, used, SKIP)


def _fill_body(slide, title: str, bullets: List[str]) -> None:
    T, S, B, SKIP = _ph_types()
    used = set()
    title_ph = next((p for p in slide.placeholders if _is_type(p, T)), None)
    if title_ph is not None:
        _set_text(title_ph, _wrap(title), _TITLE_PT, bold=True)
        used.add(title_ph.placeholder_format.idx)
    # 본문: 가장 큰 BODY/OBJECT placeholder
    bodies = [p for p in slide.placeholders if _is_type(p, B)]
    body_ph = max(bodies, key=_area) if bodies else None
    if body_ph is None:  # 제목 외 아무 placeholder
        body_ph = next((p for p in slide.placeholders
                        if p.placeholder_format.idx not in used and p.has_text_frame), None)
    if body_ph is not None:
        lines = [str(b) for b in (bullets or [])][:8] or [""]
        _set_text(body_ph, lines, _BODY_PT)
        used.add(body_ph.placeholder_format.idx)
    _clear_unused(slide, used, SKIP)


def _clear_unused(slide, used: set, skip_types) -> None:
    """채우지 않은 본문/제목성 placeholder 의 안내문구를 지운다(번호·날짜·바닥글은 보존)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx in used:
            continue
        if _is_type(ph, skip_types):
            continue
        _clear(ph)


def _wrap(text: str) -> List[str]:
    """문자열을 줄 목록으로(이미 줄바꿈 있으면 분리)."""
    if text is None:
        return [""]
    return [ln for ln in str(text).split("\n")] or [""]
