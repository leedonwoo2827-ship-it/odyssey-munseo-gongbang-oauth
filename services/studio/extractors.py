"""입력 문서 텍스트 추출 — docx / hwpx / xlsx / pptx / pdf / txt / md.

각 추출기는 의존성 미설치/파싱 실패에도 죽지 않고
"[추출 실패: ...]" 문자열을 돌려준다. (사용자가 더미 파일로 시험하므로 견고성 우선)
"""
from __future__ import annotations

import os
import zipfile
from typing import Callable, Dict

from . import config


def _clip(text: str) -> str:
    text = (text or "").strip()
    limit = config.MAX_CHARS_PER_INPUT
    if len(text) > limit:
        return text[:limit] + f"\n…(이하 {len(text) - limit}자 생략)…"
    return text


# ── 포맷별 추출기 ────────────────────────────────────────────────────
def _from_txt(path: str) -> str:
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def _from_docx(path: str) -> str:
    from docx import Document  # python-docx
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    # 표 셀 텍스트도 포함
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"# 시트: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = ["" if v is None else str(v) for v in row]
            if any(v.strip() for v in vals):
                out.append(" | ".join(vals))
    wb.close()
    return "\n".join(out)


def _from_pptx(path: str) -> str:
    from pptx import Presentation  # python-pptx
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"# 슬라이드 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        out.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
    return "\n".join(out)


def _from_pdf(path: str) -> str:
    # 1순위 pdfplumber(레이아웃 우수), 폴백 pypdf
    try:
        import pdfplumber
        out = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                out.append(page.extract_text() or "")
        text = "\n".join(out).strip()
        if text:
            return text
    except Exception:
        pass
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join((pg.extract_text() or "") for pg in reader.pages)


def _from_hwpx(path: str) -> str:
    """HWPX = ZIP 컨테이너. Contents/section*.xml 의 <hp:t> 텍스트 수집.

    표는 셀 순서대로 이어붙이며, 가능하면 행 구분을 살린다.
    """
    from lxml import etree

    HP_T = "{http://www.hancom.co.kr/hwpml/2011/paragraph}t"
    HP_P = "{http://www.hancom.co.kr/hwpml/2011/paragraph}p"

    out_lines = []
    with zipfile.ZipFile(path) as z:
        names = sorted(
            n for n in z.namelist()
            if n.startswith("Contents/section") and n.endswith(".xml")
        )
        if not names:
            # PrvText.txt(미리보기 텍스트)라도 활용
            if "Preview/PrvText.txt" in z.namelist():
                return z.read("Preview/PrvText.txt").decode("utf-8", "replace")
        for name in names:
            data = z.read(name)
            root = etree.fromstring(data)
            # 단락 단위로 텍스트를 모아 줄바꿈
            for p in root.iter(HP_P):
                texts = [t.text for t in p.iter(HP_T) if t.text]
                line = "".join(texts).strip()
                if line:
                    out_lines.append(line)
            # 단락 태그가 없는 예외적 구조 대비
            if not out_lines:
                out_lines = [t.text for t in root.iter(HP_T) if t.text and t.text.strip()]
    return "\n".join(out_lines)


_EXTRACTORS: Dict[str, Callable[[str], str]] = {
    "txt": _from_txt,
    "md": _from_txt,
    "csv": _from_txt,
    "docx": _from_docx,
    "xlsx": _from_xlsx,
    "xlsm": _from_xlsx,
    "pptx": _from_pptx,
    "pdf": _from_pdf,
    "hwpx": _from_hwpx,
}

SUPPORTED_EXTS = sorted(_EXTRACTORS.keys())


def extract(path: str) -> str:
    """경로의 파일을 텍스트로 추출. 실패 시 진단 문자열 반환(예외 던지지 않음)."""
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        return f"[지원하지 않는 형식: .{ext} — 텍스트로 시도]\n" + _safe(_from_txt, path)
    return _safe(fn, path)


def _safe(fn: Callable[[str], str], path: str) -> str:
    try:
        return _clip(fn(path))
    except ModuleNotFoundError as e:
        return f"[추출 실패: 라이브러리 미설치 ({e.name}). setup.bat 을 다시 실행하세요. 파일: {os.path.basename(path)}]"
    except Exception as e:
        return f"[추출 실패: {type(e).__name__}: {e} — 파일: {os.path.basename(path)}]"
